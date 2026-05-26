#!/usr/bin/env python3
"""
beamer_to_pptx.py

Compiles a Beamer LaTeX file to PDF and converts each slide to a .pptx file,
embedding every slide as a full-bleed EMF/SVG vector graphic.  Videos
included via the ``movie15`` package's ``\\includemovie`` command or the
``multimedia`` package's ``\\movie`` command are extracted from the source
and inserted as native PowerPoint movie shapes on the slides where they appear.

Usage:
    python beamer_to_pptx.py presentation.tex
    python beamer_to_pptx.py presentation.tex --output out.pptx
    python beamer_to_pptx.py presentation.tex --latex-engine xelatex

Requirements:
    pip install python-pptx pymupdf
    System: pdflatex (or xelatex/lualatex).
            biber and/or bibtex if the document has a bibliography.
"""

import argparse
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from collections import defaultdict
from pathlib import Path

import fitz  # PyMuPDF
import lxml.etree as etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

try:
    from pptx.opc.part import Part as _Part  # pptx < 1.0
    from pptx.opc.packuri import PackURI

    def Part(partname, content_type, blob, package):
        # pptx < 1.0: Part(partname, content_type, blob, package=None)
        return _Part(partname, content_type, blob, package)
except ImportError:
    from pptx.opc.package import Part as _Part, PackURI  # pptx >= 1.0

    def Part(partname, content_type, blob, package):
        # pptx >= 1.0: Part(partname, content_type, package, blob=b"")
        return _Part(partname, content_type, package, blob)


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────

def run(cmd: list[str], cwd: Path | None = None, check: bool = True,
        env: dict | None = None) -> subprocess.CompletedProcess:
    """Run a subprocess command, streaming output to the terminal."""
    print(f"  $ {' '.join(str(c) for c in cmd)}")
    return subprocess.run(cmd, cwd=cwd, check=check, env=env,
                          stdout=sys.stdout, stderr=sys.stderr)


def require(tool: str) -> str:
    """Return full path of *tool* or abort with a helpful message."""
    path = shutil.which(tool)
    if not path:
        sys.exit(f"[error] '{tool}' not found on PATH. Please install it and retry.")
    return path


# ──────────────────────────────────────────────────────────────────────────────
# Step 1 – Compile .tex → .pdf
# ──────────────────────────────────────────────────────────────────────────────

def compile_latex(tex_path: Path, engine: str, build_dir: Path,
                  latex_cwd: Path) -> Path:
    """
    Run the LaTeX engine (plus biber/bibtex when a bibliography is present)
    with the working directory set to *latex_cwd* so relative ``\\input``,
    ``\\includegraphics``, and ``\\bibliography`` paths resolve correctly.
    Output artefacts are written to *build_dir*.  Returns the path to the PDF.
    """
    engine_bin = require(engine)
    tex_abs = tex_path.resolve()
    stem = tex_abs.stem

    compile_cmd = [
        engine_bin,
        "-interaction=nonstopmode",
        "-halt-on-error",
        f"-output-directory={build_dir}",
        str(tex_abs),
    ]

    print(f"  LaTeX working directory: {latex_cwd}")
    print("\n[1/3] Compiling LaTeX (pass 1) …")
    run(compile_cmd, cwd=latex_cwd)

    _process_bibliography(build_dir, stem, latex_cwd)

    print("\n[1/3] Compiling LaTeX (pass 2) …")
    run(compile_cmd, cwd=latex_cwd)
    print("\n[1/3] Compiling LaTeX (pass 3 – cross-refs) …")
    run(compile_cmd, cwd=latex_cwd)

    pdf_path = build_dir / (stem + ".pdf")
    if not pdf_path.exists():
        sys.exit(f"[error] Expected PDF not found: {pdf_path}")
    print(f"  → PDF: {pdf_path}")
    return pdf_path


def _process_bibliography(build_dir: Path, stem: str, latex_cwd: Path) -> None:
    """
    Run biber or bibtex if the first LaTeX pass produced bibliography artefacts.

      *.bcf present  → biblatex + biber backend
      \\bibdata in .aux → traditional bibtex / natbib
      neither        → no bibliography, nothing to do

    Both tools are invoked with cwd=*build_dir* (where the .aux/.bcf live)
    and just the stem name as their target.  Running in *build_dir* is
    required for bibtex — TeX Live's default ``openout_any = p`` (paranoid)
    refuses to write .blg/.bbl to absolute paths outside the working
    directory.  ``BIBINPUTS``/``BSTINPUTS`` are extended to also search
    *latex_cwd* so the user's .bib files (referenced from the source by
    relative path) are still found, plus *build_dir* itself so that
    biblatex's auto-generated ``<stem>-blx.bib`` (next to the .aux) is
    picked up.  The trailing empty path element preserves the default TeX
    search locations.
    """
    bcf_path = build_dir / f"{stem}.bcf"
    aux_path = build_dir / f"{stem}.aux"

    search = os.pathsep.join([str(build_dir), str(latex_cwd), ""])
    env = {**os.environ, "BIBINPUTS": search, "BSTINPUTS": search}

    if bcf_path.exists():
        biber_bin = require("biber")
        print("\n[1/3] Processing bibliography (biber) …")
        run([biber_bin, stem], cwd=build_dir, env=env)
        return

    if aux_path.exists() and "\\bibdata" in aux_path.read_text(errors="ignore"):
        bibtex_bin = require("bibtex")
        print("\n[1/3] Processing bibliography (bibtex) …")
        run([bibtex_bin, stem], cwd=build_dir, env=env)


# ──────────────────────────────────────────────────────────────────────────────
# Step 2 – Split PDF into per-slide SVG vector files
# ──────────────────────────────────────────────────────────────────────────────

def pdf_to_svgs_pymupdf(pdf_path: Path, svg_dir: Path) -> list[Path]:
    """
    Use PyMuPDF (fitz) to render each PDF page as an SVG.
    PyMuPDF produces true vector SVGs (text + paths), not bitmaps.
    """
    svg_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    paths = []
    for i, page in enumerate(doc):
        svg_bytes = page.get_svg_image(matrix=fitz.Identity)
        out = svg_dir / f"slide_{i+1:04d}.svg"
        out.write_bytes(svg_bytes if isinstance(svg_bytes, bytes)
                        else svg_bytes.encode())
        paths.append(out)
    doc.close()
    print(f"  → {len(paths)} SVG(s) written to {svg_dir}")
    return paths


def extract_svgs(pdf_path: Path, svg_dir: Path) -> list[Path]:
    """Convert PDF pages to SVG using PyMuPDF."""
    print("\n[2/3] Converting PDF pages to SVG …")

    paths = pdf_to_svgs_pymupdf(pdf_path, svg_dir)
    if paths:
        return paths

    sys.exit(
        "[error] PyMuPDF is required for SVG extraction.\n"
        "Install it with:\n"
        "  pip install pymupdf"
    )


# ──────────────────────────────────────────────────────────────────────────────
# Video discovery — supports both the ``movie15`` package's ``\includemovie``
# and the ``multimedia`` package's ``\movie`` command.
# ──────────────────────────────────────────────────────────────────────────────

# movie15: \includemovie[opts]{width}{height}{path}
_INCLUDEMOVIE_RE = re.compile(
    r'\\includemovie\s*'
    r'(?:\[(?P<opts>[^\]]*)\])?\s*'
    r'\{(?P<w>[^}]*)\}\s*'
    r'\{(?P<h>[^}]*)\}\s*'
    r'\{(?P<path>[^}]*)\}',
    re.DOTALL,
)

# multimedia: \movie[opts]{poster}{path}
# Cannot use a single regex because the poster argument is often
# ``\includegraphics{…}`` — i.e. it contains balanced braces.  We just locate
# the command and parse the arguments separately with a brace-aware scanner.
_MOVIE_CMD_RE = re.compile(r'\\movie\b')

# Looks for ``\includegraphics[opts]{file}`` inside the poster argument so
# we can hand the same image to PowerPoint as the video's poster frame.
_INCLUDEGRAPHICS_RE = re.compile(
    r'\\includegraphics\s*(?:\[[^\]]*\])?\s*\{([^}]*)\}'
)


def _parse_balanced_braces(text: str, start: int) -> tuple[str, int] | None:
    """
    *start* must index a ``{``.  Return the substring inside the matching
    pair and the index just past the closing ``}``, or ``None`` on mismatch.
    Honours TeX-style backslash escapes so that ``\\{`` / ``\\}`` don't
    perturb the depth counter.
    """
    if start >= len(text) or text[start] != '{':
        return None
    depth = 1
    i = start + 1
    while i < len(text):
        ch = text[i]
        if ch == '\\' and i + 1 < len(text):
            i += 2  # skip the escaped character
            continue
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                return text[start + 1:i], i + 1
        i += 1
    return None


def _video_options(opts: str) -> dict:
    """
    Translate the comma-separated option string from ``\\includemovie`` /
    ``\\movie`` into the keys we care about: ``autoplay``, ``loop``,
    ``volume``, ``poster_path``.

    Option spellings handled:
      autoplay (movie15)      → autoplay
      autostart (multimedia)  → autoplay
      repeat[=N] (movie15)    → loop
      loop      (multimedia)  → loop
      palindrome (both)       → loop
      volume=N  (movie15)     → volume (0.0–1.0)
      poster=FILE (movie15)   → poster_path
    """
    poster_match = re.search(r'\bposter\s*=\s*([^,\s\]]+)', opts)
    volume_match = re.search(r'\bvolume\s*=\s*([0-9.]+)', opts)
    return {
        "poster_path": (poster_match.group(1).strip()
                        if poster_match else None),
        "autoplay": bool(re.search(r'\b(?:autoplay|autostart)\b', opts)),
        "loop": bool(re.search(r'\b(?:repeat|loop)\b(?:\s*=\s*\d+)?', opts)
                     or re.search(r'\bpalindrome\b', opts)),
        "volume": (float(volume_match.group(1))
                   if volume_match else None),
    }


def _parse_tex_videos(tex_path: Path) -> list[dict]:
    """
    Find every video inclusion in *tex_path* and return one dict per call,
    in source order, with keys:

      * ``video_path``  (str) — raw path as written in the source
      * ``poster_path`` (str | None) — poster image if one was specified
      * ``autoplay``    (bool)
      * ``loop``        (bool)
      * ``volume``      (float | None) — 0.0–1.0
    """
    try:
        text = tex_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = tex_path.read_text(encoding="latin-1")

    # Simple comment stripping (naive but handles % on its own line or after cmd)
    # We only strip for scanning purposes to find [ and { correctly.
    def skip_noise(t: str, p: int) -> int:
        while p < len(t):
            if t[p].isspace():
                p += 1
            elif t[p] == '%':
                nl = t.find('\n', p)
                if nl == -1: return len(t)
                p = nl + 1
            else:
                break
        return p

    found = []  # list of (source_pos, info_dict)

    # ── movie15 ──────────────────────────────────────────────────────────────
    for m in _INCLUDEMOVIE_RE.finditer(text):
        info = _video_options(m.group("opts") or "")
        info["video_path"] = m.group("path").strip()
        print(f"  [debug] matched movie15: {info['video_path']}")
        found.append((m.start(), info))

    # ── multimedia (\movie) ──────────────────────────────────────────────────
    for m in _MOVIE_CMD_RE.finditer(text):
        pos = m.end()
        # Optional [opts]
        opts = ""
        pos = skip_noise(text, pos)
        if pos < len(text) and text[pos] == '[':
            close = text.find(']', pos)
            if close == -1:
                print(f"  [debug] \movie at {m.start()} has unclosed [")
                continue
            opts = text[pos + 1:close]
            pos = close + 1
        
        pos = skip_noise(text, pos)
        if pos >= len(text) or text[pos] != '{':
            print(f"  [debug] \movie at {m.start()} missing poster arg at {pos}: {text[pos:pos+10]!r}")
            continue
        first = _parse_balanced_braces(text, pos)
        if first is None:
            print(f"  [debug] \movie at {m.start()} failed to parse poster arg braces")
            continue
        poster_arg, pos = first
        
        pos = skip_noise(text, pos)
        if pos >= len(text) or text[pos] != '{':
            print(f"  [debug] \movie at {m.start()} missing path arg at {pos}: {text[pos:pos+10]!r}")
            continue
        second = _parse_balanced_braces(text, pos)
        if second is None:
            print(f"  [debug] \movie at {m.start()} failed to parse path arg braces")
            continue
        path_arg, _ = second

        info = _video_options(opts)
        info["video_path"] = path_arg.strip()
        print(f"  [debug] matched multimedia at {m.start()}: {info['video_path']}")
        # If the poster argument uses \includegraphics, treat that image as
        # the poster.  Otherwise (plain text, \hyperlink, etc.) leave it as
        # whatever the explicit poster= option said, which may be None.
        if info["poster_path"] is None:
            ig = _INCLUDEGRAPHICS_RE.search(poster_arg)
            if ig:
                info["poster_path"] = ig.group(1).strip()
        found.append((m.start(), info))

    found.sort(key=lambda x: x[0])
    return [info for _, info in found]


def _find_video_annotations(pdf_path: Path) -> list[tuple[int, "fitz.Rect"]]:
    """
    Inspect the compiled PDF for Screen / Movie annotations.
    """
    doc = fitz.open(str(pdf_path))
    results = []
    try:
        for page_idx in range(doc.page_count):
            page = doc[page_idx]
            annots = page.annots()
            if annots is None:
                continue
            for annot in annots:
                _, name = annot.type
                print(f"  [debug] page {page_idx} annot type: {name}")
                # multimedia sometimes uses Screen, sometimes Movie
                if name in ("Screen", "Movie"):
                    results.append((page_idx, annot.rect))
    finally:
        doc.close()
    return results


def _collect_videos(tex_path: Path, pdf_path: Path,
                    source_dir: Path) -> dict[int, list[dict]]:
    """
    Combine .tex parsing and PDF annotation analysis.  Returns a mapping
    ``{page_index_0based: [video_info, ...]}`` where ``video_info`` has resolved
    absolute ``video_path`` / ``poster_path`` and a PyMuPDF
    ``rect_pt`` describing the on-page position in points.

    If the number of ``\\includemovie`` / ``\\movie`` calls in the source
    doesn't match the number of Screen/Movie annotations in the PDF, we
    fall back to the smaller count and pair them in document order (extra
    entries on either side are dropped with a warning).
    """
    tex_videos = _parse_tex_videos(tex_path)
    if not tex_videos:
        return {}

    annots = _find_video_annotations(pdf_path)
    if not annots:
        print("  [warn] video commands found in source but no Screen/Movie "
              "annotations in the PDF — videos cannot be positioned and "
              "will be skipped.")
        return {}

    if len(annots) != len(tex_videos):
        print(f"  [warn] {len(tex_videos)} video command(s) in source "
              f"but {len(annots)} video annotation(s) in PDF — pairing the "
              f"first {min(len(annots), len(tex_videos))} in document order.")

    result = defaultdict(list)
    for video, (page_idx, rect) in zip(tex_videos, annots):
        video_path = (source_dir / video["video_path"]).resolve()
        poster_path = None
        if video["poster_path"]:
            poster_path = (source_dir / video["poster_path"]).resolve()
        result[page_idx].append({
            **video,
            "video_path": video_path,
            "poster_path": poster_path,
            "rect_pt": rect,
        })
    return result


def _add_video_to_slide(slide, info: dict) -> None:
    """Insert a native PowerPoint movie shape at the PDF annotation's rect."""
    video_path: Path = info["video_path"]
    poster_path: Path | None = info["poster_path"]
    rect = info["rect_pt"]

    if not video_path.exists():
        print(f"  [warn] video file not found, skipping: {video_path}")
        return

    # PyMuPDF Rect uses points with the origin at the top-left of the page;
    # PowerPoint EMU uses the same origin convention.  1 pt = 12 700 EMU.
    left = int(rect.x0 * 12700)
    top = int(rect.y0 * 12700)
    width = int(rect.width * 12700)
    height = int(rect.height * 12700)

    kwargs = {}
    if poster_path and poster_path.exists():
        kwargs["poster_frame_image"] = str(poster_path)

    shape = slide.shapes.add_movie(
        str(video_path), left, top, width, height, **kwargs,
    )

    autoplay = info.get("autoplay", False)
    loop = info.get("loop", False)
    volume = info.get("volume")  # 0.0..1.0 or None
    if autoplay or loop or volume is not None:
        _apply_video_playback(slide, shape,
                              autoplay=autoplay, loop=loop, volume=volume)

    extras = []
    if autoplay: extras.append("autoplay")
    if loop:     extras.append("loop")
    if volume is not None: extras.append(f"vol={volume:.2f}")
    extras_str = f" [{', '.join(extras)}]" if extras else ""
    print(f"  → embedded video {video_path.name} on slide "
          f"(at {rect.x0:.0f},{rect.y0:.0f} pt, "
          f"{rect.width:.0f}×{rect.height:.0f} pt){extras_str}")


# Full <p:timing> tree that triggers ``playFrom(0.0)`` on the named shape as
# soon as the slide is shown.  The ``presetClass="mediacall"`` afterEffect is
# the structure PowerPoint itself emits for autoplay videos.  ``{shape_id}``,
# ``{vol}`` and ``{repeat_attr}`` are substituted by Python's ``str.format``.
_AUTOPLAY_TIMING_XML = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="indefinite"/>
                                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:video>
            <p:cMediaNode vol="{vol}">
              <p:cTn id="7" fill="hold" display="0"{repeat_attr}>
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


# Click-to-play timing — same shape as what python-pptx emits by default but
# with the ``vol`` attribute and an optional ``repeatCount`` so we can honour
# ``volume=`` / ``repeat`` even when the user did not request autoplay.
_CLICK_TIMING_XML = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:video>
            <p:cMediaNode vol="{vol}">
              <p:cTn id="2" fill="hold" display="0"{repeat_attr}>
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


def _apply_video_playback(slide, shape, *,
                          autoplay: bool, loop: bool,
                          volume: float | None) -> None:
    """
    Replace the slide's <p:timing> with a tree that wires *shape* for
    autoplay, looping, and/or custom volume.  python-pptx's ``add_movie``
    always leaves a click-to-play <p:timing> behind, so we just swap it.
    """
    shape_id = shape.shape_id
    # OOXML ``vol`` is on the 0..100000 (= 100 %) scale; the PowerPoint
    # default for unspecified volume is 80 %.
    vol = int((volume if volume is not None else 0.8) * 100000)
    vol = max(0, min(100000, vol))
    repeat_attr = ' repeatCount="indefinite"' if loop else ''

    template = _AUTOPLAY_TIMING_XML if autoplay else _CLICK_TIMING_XML
    timing_xml = template.format(shape_id=shape_id, vol=vol,
                                 repeat_attr=repeat_attr)
    new_timing = etree.fromstring(timing_xml)

    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    sld = slide.element
    existing = sld.find(f"{P}timing")
    if existing is not None:
        sld.remove(existing)
    sld.append(new_timing)


# ──────────────────────────────────────────────────────────────────────────────
# Step 3 – Build .pptx
# ──────────────────────────────────────────────────────────────────────────────

def _svg_dimensions(svg_path: Path) -> tuple[float, float]:
    """
    Parse width/height from the <svg> root element (in pt).
    Falls back to 4:3 Beamer default (362.835 × 272.126 pt ≈ 128×96 mm).
    """
    BEAMER_W_PT = 362.835   # 128 mm
    BEAMER_H_PT = 272.126   # 96 mm

    try:
        tree = ET.parse(str(svg_path))
        root = tree.getroot()
        # Strip namespace if present
        attrib = root.attrib

        def to_pt(value: str) -> float:
            value = value.strip()
            if value.endswith("pt"):
                return float(value[:-2])
            if value.endswith("mm"):
                return float(value[:-2]) * 2.8346456693
            if value.endswith("cm"):
                return float(value[:-2]) * 28.346456693
            if value.endswith("in"):
                return float(value[:-2]) * 72.0
            if value.endswith("px"):
                return float(value[:-2]) * 0.75  # 96 dpi → pt
            return float(value)  # assume pt

        w = to_pt(attrib["width"])  if "width"  in attrib else BEAMER_W_PT
        h = to_pt(attrib["height"]) if "height" in attrib else BEAMER_H_PT
        return w, h

    except Exception:
        return BEAMER_W_PT, BEAMER_H_PT


def build_pptx(svg_paths: list[Path], output_path: Path,
               videos: dict[int, list[dict]] | None = None) -> None:
    """
    Create a .pptx where each slide contains one full-bleed SVG picture.

    If *videos* is supplied (mapping ``page_index_0based → [video_info, ...]``),
    native PowerPoint movie shapes are added on top of the SVG.
    """
    print("\n[3/3] Building PPTX …")

    videos = videos or {}

    # Infer slide dimensions from the first SVG
    slide_w_pt, slide_h_pt = _svg_dimensions(svg_paths[0])
    print(f"  Slide dimensions: {slide_w_pt:.2f} × {slide_h_pt:.2f} pt")

    prs = Presentation()

    # Set slide size to match the source PDF pages (EMU = pt × 12700)
    prs.slide_width  = int(slide_w_pt * 12700)
    prs.slide_height = int(slide_h_pt * 12700)

    # Use a completely blank layout (no placeholders)
    blank_layout = prs.slide_layouts[6]

    for idx, svg_path in enumerate(svg_paths, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # python-pptx does not natively support SVG insertion into the
        # slide XML, so we add it via direct XML manipulation (Office 2016+
        # supports <a:graphicData uri="…svg…"> inside a picture shape).
        # For maximum compatibility we also embed a PNG fallback.
        _insert_svg_picture(slide, svg_path,
                            prs.slide_width, prs.slide_height)

        # Overlay any videos on top of the SVG.
        page_idx = idx - 1
        if page_idx in videos:
            for video_info in videos[page_idx]:
                _add_video_to_slide(slide, video_info)

        if idx % 10 == 0 or idx == len(svg_paths):
            print(f"  Processed {idx}/{len(svg_paths)} slides …")

    prs.save(str(output_path))
    print(f"\n✓ Saved: {output_path}")


# ──────────────────────────────────────────────────────────────────────────────
# SVG embedding helpers
# ──────────────────────────────────────────────────────────────────────────────

def _png_fallback(svg_path: Path) -> bytes:
    """
    Render SVG → PNG bytes.

    IMPORTANT: <a:blip r:embed="..."> pointing at a raster image is *required*
    by the OOXML spec — the <asvg:svgBlip> is only an extension on top of it.
    PowerPoint silently ignores the SVG and falls back to rasterising (or shows
    nothing) when <a:blip> has no r:embed.  This function therefore always
    returns valid PNG bytes; it aborts the process if every renderer fails.

    Uses PyMuPDF (fitz) to rasterise the SVG.
    """
    try:
        doc = fitz.open(stream=svg_path.read_bytes(), filetype="svg")
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(1, 1))
        png = pix.tobytes("png")
        doc.close()
        return png
    except Exception as exc:
        sys.exit(
            f"[error] Cannot render SVG → PNG via PyMuPDF: {exc}\n"
            "A raster fallback is required for PowerPoint to honour the "
            "embedded SVG."
        )


def _insert_svg_picture(slide, svg_path: Path, cx_emu: int, cy_emu: int) -> None:
    """
    Embed *svg_path* as a full-slide picture using the Office Open XML SVG
    extension (DrawingML / OOXML §20.1.8.55).

    OOXML structure required by PowerPoint to actually honour the SVG:

      <a:blip r:embed="PNG_rId">              ← MANDATORY raster base image
        <a:extLst>
          <a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">
            <asvg:svgBlip r:embed="SVG_rId"/> ← vector overlay (Office 2016+)
          </a:ext>
        </a:extLst>
      </a:blip>

    The ``{96DAC541-…}`` URI is the well-known DrawingML extension identifier
    for SVG-aware blips.  Without the <a:extLst>/<a:ext> wrapper PowerPoint
    silently drops the <asvg:svgBlip> and renders only the PNG fallback.
    """
    svg_bytes = svg_path.read_bytes()
    # Always produces bytes or aborts — see _png_fallback docstring.
    png_bytes = _png_fallback(svg_path)

    slide_part = slide.part
    package = slide_part.package

    # PNG part — registered as a normal image relationship (r:embed on <a:blip>)
    try:
        _png_part, png_rId = slide_part.get_or_add_image_part(io.BytesIO(png_bytes))
    except Exception:
        png_partname = PackURI(f"/ppt/media/{svg_path.stem}_{id(svg_path)}_fb.png")
        png_part = Part(png_partname, "image/png", png_bytes, package)
        png_rId = slide_part.relate_to(png_part, RT.IMAGE)

    # SVG part — must use the *standard* image relationship type, just like the
    # PNG.  PowerPoint-generated pptx files always relate SVG attachments via
    # RT.IMAGE; the bespoke ``…/relationships/svgBlip`` URI is not recognised
    # by the renderer, which then refuses to load the picture entirely.
    # The ``package=`` backref is required so the part is correctly serialised
    # into the .pptx archive (without it, PowerPoint's resolver fails and the
    # whole picture shape collapses into a “Picture can't be displayed” box).
    svg_partname = PackURI(f"/ppt/media/{svg_path.stem}_{id(svg_path)}.svg")
    svg_part = Part(svg_partname, "image/svg+xml", svg_bytes, package)
    svg_rId = slide_part.relate_to(svg_part, RT.IMAGE)

    # ── Build the whole <p:pic> as one piece of XML so lxml uses the prefixes
    # declared on the root element (otherwise it invents ns0/ns1/… prefixes on
    # the <a:blip> subtree, which some readers refuse to load).
    SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
    pic_xml = f"""\
<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main">
  <p:nvPicPr>
    <p:cNvPr id="2" name="Slide {svg_path.stem}"/>
    <p:cNvPicPr preferRelativeResize="0"/>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{png_rId}">
      <a:extLst>
        <a:ext uri="{SVG_EXT_URI}">
          <asvg:svgBlip r:embed="{svg_rId}"/>
        </a:ext>
      </a:extLst>
    </a:blip>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="0" y="0"/>
      <a:ext cx="{cx_emu}" cy="{cy_emu}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>
"""
    pic_elem = etree.fromstring(pic_xml)

    # Append <p:pic> to the slide shape tree.
    slide.shapes._spTree.append(pic_elem)


# ──────────────────────────────────────────────────────────────────────────────
# CLI
# ──────────────────────────────────────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Compile a Beamer .tex file to PDF and export slides as "
                    "vector SVG graphics inside a .pptx presentation. "
                    "A compiled .pdf file can also be used directly."
    )
    p.add_argument("input_file", type=Path,
                   help="Path to the Beamer .tex source file or an existing .pdf file")
    p.add_argument("--output", "-o", type=Path, default=None,
                   help="Output .pptx path (default: <input_stem>.pptx alongside the source)")
    p.add_argument("--latex-engine", "-e", default="pdflatex",
                   choices=["pdflatex", "xelatex", "lualatex"],
                   help="LaTeX engine to use (default: pdflatex)")
    p.add_argument("--keep-build", action="store_true",
                   help="Keep the temporary build directory after completion "
                        "(ignored when --build-dir is set, as that dir is never deleted)")
    p.add_argument("--build-dir", "-b", type=Path, default=None,
                   metavar="DIR",
                   help="Use DIR as the build directory instead of a system temp dir. "
                        "The directory is created if it does not exist and is never "
                        "deleted automatically (implies --keep-build).")
    p.add_argument("--latex-cwd", type=Path, default=None,
                   metavar="DIR",
                   help="Working directory for LaTeX/biber/bibtex "
                        "(default: directory containing the .tex file). "
                        "Affects how relative \\input, \\includegraphics, "
                        "and \\bibliography paths in the source resolve.")
    return p.parse_args()


def _run_build(args: argparse.Namespace, build_dir: Path) -> None:
    """Core build logic, shared by the temp-dir and user-defined-dir paths."""
    input_path: Path = args.input_file.resolve()
    output_path: Path = args.output or input_path.with_suffix(".pptx")
    svg_dir = build_dir / "svgs"
    latex_cwd = (args.latex_cwd or input_path.parent).resolve()

    # ── Compile or use existing PDF ───────────────────────────────────────────
    if input_path.suffix.lower() == ".pdf":
        # Input is already a PDF
        pdf_path = input_path
        print(f"[1/3] Using input PDF: {pdf_path}")
    else:
        # Input is assumed to be LaTeX
        pdf_path = compile_latex(input_path, args.latex_engine, build_dir, latex_cwd)

    # ── Extract SVGs ──────────────────────────────────────────────────────────
    svg_paths = extract_svgs(pdf_path, svg_dir)

    # ── Discover movie15 videos ──────────────────────────────────────────────
    # Only try to collect videos if we have a .tex file
    videos = {}
    if input_path.suffix.lower() == ".tex":
        videos = _collect_videos(input_path, pdf_path, latex_cwd)
        total_videos = sum(len(vlist) for vlist in videos.values())
        if total_videos:
            print(f"  Detected {total_videos} video(s) to embed.")

    # ── Build PPTX ────────────────────────────────────────────────────────────
    build_pptx(svg_paths, output_path, videos=videos)


def main() -> None:
    args = parse_args()

    input_path: Path = args.input_file.resolve()
    if not input_path.exists():
        sys.exit(f"[error] File not found: {input_path}")

    if args.build_dir:
        # User-defined build directory — create it and never auto-delete it.
        build_dir = args.build_dir.resolve()
        build_dir.mkdir(parents=True, exist_ok=True)
        print(f"  Build directory: {build_dir}")
        _run_build(args, build_dir)
    else:
        # Default: auto-managed temp directory.
        with tempfile.TemporaryDirectory(prefix="beamer2pptx_") as tmp:
            build_dir = Path(tmp)
            _run_build(args, build_dir)

            if args.keep_build:
                kept = input_path.parent / (input_path.stem + "_build")
                shutil.copytree(tmp, str(kept), dirs_exist_ok=True)
                print(f"  Build directory kept at: {kept}")

    print("\nDone.")


if __name__ == "__main__":
    main()
